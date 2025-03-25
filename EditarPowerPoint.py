from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE 
import os
import win32com.client
from CrearImagenes import powerpoint_to_images




# Definir las dimensiones específicas para eliminar la imagen
dimensiones_deseadas = {
    'width': 708571,  # Reemplaza con el valor deseado en EMUs
    'height': 656136  # Reemplaza con el valor deseado en EMUs
}


def quitarElementos(archivo_pptx):

    eliminar_diapositivas(os.path.abspath(archivo_pptx))

    # Cargar la presentación
    if archivo_pptx:
        presentation = Presentation(archivo_pptx)
        # Palabras clave a buscar
        palabras_clave = ["Resultado", "Mapa", "Siguiente", "Observaciones", "Información", "Anterior"]

        # Recorrer las diapositivas
        for slide in presentation.slides:
            # Identificar cuadros de texto para eliminar
            for shape in slide.shapes:

                if shape.shape_type == 13:  # 13 indica una imagen
                    # Obtener las dimensiones de la imagen
                    width = shape.width
                    height = shape.height
                    
                    # Verificar si las dimensiones coinciden con las deseadas
                    if width == dimensiones_deseadas['width'] and height == dimensiones_deseadas['height']:
                        # Eliminar la imagen
                        sp = shape
                        sp._element.getparent().remove(sp._element)  
                
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    if shape.has_text_frame:   
                        texto = shape.text_frame.text.strip()
                        if any(palabra in texto for palabra in palabras_clave):
                            sp = shape
                            sp._element.getparent().remove(sp._element) 
                
        # Guardar los cambios en el mismo archivo
        presentation.save(archivo_pptx)
        print("Tablas y elementos eliminados. Archivo guardado.")
    else:
        print("No se seleccionó ningún archivo.")

def eliminar_diapositivas(archivo_pptx):
    # Crear instancia de PowerPoint
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    
    # Abrir la presentación en modo de solo lectura
    presentation = powerpoint.Presentations.Open(archivo_pptx, WithWindow=False)

    #Total de diapositivas
    totalDiapositivas = presentation.Slides.count

    #Eliminar ultima diapositiva
    presentation.Slides(totalDiapositivas).Delete()
    
    #Eliminar primeras 4 diapostivas
    diapositivas_a_eliminar = [1, 2, 3, 4]
    for indice in sorted(diapositivas_a_eliminar, reverse=True):
        presentation.Slides(indice).Delete()

    # Guarda la presentación
    presentation.SaveAs(archivo_pptx)

    # Cierra la presentación
    presentation.Close()

    # Cierra PowerPoint
    powerpoint.Quit()


